"""Generates a PowerPoint presentation (.pptx) summarizing the AI IT Support
Assistant project - for demonstration/evaluation purposes.

Dev-only script (python-pptx is not an app dependency, only needed here).

Run:  python scripts/generate_ppt.py
Output: docs/AI_IT_Support_Assistant_Presentation.pptx
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "docs" / "AI_IT_Support_Assistant_Presentation.pptx"

NAVY = RGBColor(0x0F, 0x2A, 0x43)
BLUE = RGBColor(0x1F, 0x6F, 0xEB)
LIGHT_BLUE = RGBColor(0xE8, 0xF1, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
GREEN = RGBColor(0x2E, 0xA0, 0x4A)
AMBER = RGBColor(0xC9, 0x7A, 0x0A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=NAVY, align=PP_ALIGN.LEFT, font="Segoe UI"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=20, color=NAVY, line_spacing=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return box


def add_footer(slide, page_no):
    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(6), Inches(0.3),
                "AI IT Support Assistant - Agentic AI Final Project", size=10, color=GRAY)
    add_textbox(slide, Inches(12.6), Inches(7.1), Inches(0.5), Inches(0.3),
                str(page_no), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Segoe UI"


def content_slide(prs, title, bullets, page_no, size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_header_bar(slide, title)
    add_bullets(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.2), bullets, size=size)
    add_footer(slide, page_no)
    return slide


def box(slide, left, top, width, height, text, fill=BLUE, font_color=WHITE, size=14, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = font_color
    run.font.name = "Segoe UI"
    return shp


def arrow_down(slide, left, top, width=Inches(0.5), height=Inches(0.35)):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GRAY
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    page = 1

    # --- Slide 1: Title -----------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_textbox(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2),
                "AI IT Support Assistant", size=48, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.8),
                "Agentic AI Helpdesk Assistant  -  LangChain + LangGraph + Streamlit",
                size=22, color=RGBColor(0xBF, 0xD6, 0xF7))
    add_textbox(slide, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.5),
                "GenAI Development Program - Final Evaluation Project (Project 3: Agentic AI)",
                size=16, color=RGBColor(0x8F, 0xAE, 0xD6))
    add_textbox(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
                "No external database. Runs fully local. LLM optional.", size=14, color=RGBColor(0x8F, 0xAE, 0xD6))
    page += 1

    # --- Slide 2: Problem Statement ------------------------------------------
    content_slide(prs, "Problem Statement", [
        "Employees need fast, self-service help for common IT issues",
        "Manual channels are slow: searching a wiki, emailing IT, opening a portal ticket",
        "No single place to: get troubleshooting help, check ticket status, raise a ticket, or check service status",
        "Goal: an AI agent that understands the request, picks the right tool, validates before acting, and never invents information",
    ], page)
    page += 1

    # --- Slide 3: Solution Overview -------------------------------------------
    content_slide(prs, "Solution Overview", [
        "A LangGraph state machine wrapped in a Streamlit chat UI",
        "Intent/Decision node decides if a tool is needed and which one",
        "Conditional edges route to the correct tool node (knowledge search, ticket lookup, ticket creation, employee lookup, system status)",
        "Response Generation node phrases the tool result into a clear, safety-checked answer",
        "Conversation memory persists per session via LangGraph's MemorySaver checkpointer",
        "Works fully offline with a deterministic rule-based fallback when no LLM API key is configured",
    ], page)
    page += 1

    # --- Slide 4: System Architecture (diagram) -------------------------------
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_header_bar(slide, "System Architecture")
    layers = [
        ("Streamlit Chat UI  (app.py)", BLUE),
        ("LangGraph Orchestrator  +  MemorySaver Checkpointer", NAVY),
        ("Reasoning Layer:  LLM tool-calling (OpenAI / Gemini)  OR  Rule-based fallback", AMBER),
        ("Tool Layer:  knowledge_search | ticket_lookup | create_ticket | employee_lookup | system_status_check", GREEN),
        ("Data Layer:  In-memory SQLite  (seeded from data/*.json)", GRAY),
    ]
    top = Inches(1.5)
    box_h = Inches(0.85)
    gap = Inches(0.25)
    for label, color in layers:
        box(slide, Inches(1.2), top, Inches(10.9), box_h, label, fill=color, size=16)
        top = Emu(top + box_h)
        if label != layers[-1][0]:
            arrow_down(slide, Inches(6.35), top, width=Inches(0.6), height=Emu(gap))
        top = Emu(top + gap)
    add_footer(slide, page)
    page += 1

    # --- Slide 5: LangGraph Workflow (diagram) --------------------------------
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_header_bar(slide, "LangGraph Workflow")
    box(slide, Inches(5.15), Inches(1.5), Inches(3.0), Inches(0.7), "classify_intent\n(Intent / Decision Node)", fill=NAVY, size=13)
    branches = [
        ("employee_lookup", GREEN),
        ("knowledge_search", GREEN),
        ("ticket_lookup", GREEN),
        ("ticket_creation\n(validate + confirm)", AMBER),
        ("system_status", GREEN),
        ("direct_response\n(chit-chat / cancel)", GRAY),
    ]
    n = len(branches)
    total_w = Inches(12.2)
    box_w = Inches(1.85)
    start_x = Inches(0.55)
    top_y = Inches(2.9)
    for i, (label, color) in enumerate(branches):
        x = Emu(int(start_x) + i * int((total_w - box_w) / (n - 1)))
        box(slide, x, top_y, box_w, Inches(0.9), label, fill=color, size=11)
    box(slide, Inches(3.6), Inches(4.5), Inches(3.2), Inches(0.7), "response_generation", fill=BLUE, size=14)
    box(slide, Inches(9.8), Inches(4.5), Inches(2.2), Inches(0.7), "(direct_response\nskips this node)", fill=WHITE, font_color=GRAY, size=11, shape=MSO_SHAPE.RECTANGLE)
    box(slide, Inches(4.9), Inches(5.7), Inches(3.5), Inches(0.7), "Final Answer -> Streamlit UI", fill=NAVY, size=14)
    add_textbox(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
                "Conditional routing also re-enters this flow for multi-turn confirmations (employee ID, yes/no) using state.pending_action",
                size=12, color=GRAY)
    add_footer(slide, page)
    page += 1

    # --- Slide 6: Required Tools ------------------------------------------------
    content_slide(prs, "Agent Tools", [
        "Tool 1 - Knowledge Search: searches the local IT knowledge base and answers how-to / troubleshooting questions",
        "Tool 2 - Ticket Lookup: searches existing support tickets for a verified employee and returns their status",
        "Tool 3 - Ticket Creation: validates required info, checks for duplicates, confirms with the user, then creates the ticket",
        "Bonus Tool - System Status Check: reports live operational status of IT services (VPN, Email, WiFi, Printers, HR Portal)",
        "Supporting Tool - Employee Lookup: verifies an employee ID against the company directory before any ticket action",
    ], page)
    page += 1

    # --- Slide 7: Safety & Validation --------------------------------------------
    content_slide(prs, "Safety & Validation", [
        "Never creates a ticket without an explicit user confirmation (yes/no)",
        "Duplicate-ticket detection: flags similar open tickets before creating a new one",
        "Missing information (employee ID, issue description) is always asked for, never assumed",
        "Tool failures and unknown employee IDs are handled gracefully with clear messages",
        "\"Escape hatch\": if the user changes topic mid-confirmation, the agent re-classifies instead of getting stuck",
        "Retrieved facts (tool results) vs. LLM-generated phrasing are kept clearly distinct - no invented ticket IDs, names, or KB content",
    ], page)
    page += 1

    # --- Slide 8: Memory & State --------------------------------------------------
    content_slide(prs, "Memory / State Management", [
        "AgentState (TypedDict): messages, employee identity, pending_action, ticket draft, tool trace",
        "LangGraph MemorySaver checkpointer persists state per browser session (thread_id)",
        "Example multi-turn flow:",
        "   User: \"I have a VPN issue\"  ->  Agent: \"What is your employee ID?\"",
        "   User: \"EMP1024\"  ->  Agent: \"Found your profile. Check existing tickets?\"",
        "   User: \"Yes\"  ->  Agent returns the employee's tickets",
    ], page)
    page += 1

    # --- Slide 9: Tech Stack --------------------------------------------------
    content_slide(prs, "Technology Stack", [
        "Agent orchestration: LangGraph (StateGraph, conditional edges, MemorySaver)",
        "LLM / tool calling: LangChain (bind_tools), OpenAI or Google Gemini (optional)",
        "Fallback reasoning: deterministic rule-based classifier (no API key required)",
        "Data storage: SQLite, in-memory only - no external database",
        "UI: Streamlit chat interface",
        "Config & logging: python-dotenv + typed Settings, Python logging module",
    ], page)
    page += 1

    # --- Slide 10: Project Structure --------------------------------------------
    content_slide(prs, "Project Structure", [
        "app.py - Streamlit entry point",
        "data/ - sample employees, tickets, knowledge base, system status (JSON)",
        "src/database/ - in-memory SQLite schema, seeding, reset",
        "src/tools/ - LangChain @tool functions (one per capability)",
        "src/agent/ - state.py, llm.py, rules.py, nodes.py, graph.py",
        "sample_outputs/ - captured end-to-end conversation transcript",
        "docs/ - architecture diagrams",
    ], page, size=19)
    page += 1

    # --- Slide 11: Sample Interaction --------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_header_bar(slide, "Sample Interaction")
    convo = [
        ("User", "My VPN is not working, please raise a ticket"),
        ("Assistant", "I'd like to create the following ticket:\nCategory: VPN | Priority: Medium\nDescription: My VPN is not working\nShall I proceed? (yes/no)"),
        ("User", "Yes"),
        ("Assistant", "\u2705 Your ticket has been created: TCK-1006 (VPN, Medium priority, status: Open)."),
    ]
    top = Inches(1.5)
    for speaker, text in convo:
        is_user = speaker == "User"
        fill = LIGHT_BLUE if is_user else NAVY
        font_color = NAVY if is_user else WHITE
        h = Inches(0.5 + 0.28 * (text.count("\n") + 1))
        box(slide, Inches(0.9) if is_user else Inches(4.4), top, Inches(8.0), h,
            f"{speaker}:  {text}", fill=fill, font_color=font_color, size=13,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        top = Emu(int(top) + int(h) + int(Inches(0.25)))
    add_footer(slide, page)
    page += 1

    # --- Slide 12: Key Design Decisions -------------------------------------------
    content_slide(prs, "Key Design Decisions", [
        "In-memory SQLite (no external DB) seeded from JSON - satisfies the local-only requirement",
        "Hybrid reasoning: LLM function-calling when available, deterministic rules always available as a safety net",
        "Mandatory confirmation before every ticket creation - agent never acts blindly",
        "LangGraph MemorySaver for real conversation memory, not manual session hacks",
        "Duplicate-ticket detection using keyword overlap on open/in-progress tickets",
    ], page, size=19)
    page += 1

    # --- Slide 13: Limitations --------------------------------------------------
    content_slide(prs, "Limitations", [
        "In-memory DB resets when the Streamlit process restarts (by design)",
        "Rule-based fallback understands fewer phrasings than a genuine LLM",
        "Duplicate detection is a keyword-overlap heuristic, not semantic similarity",
        "Single shared DB connection per server process (demo scope, not multi-tenant)",
        "No authentication layer - local demo assistant, not a production helpdesk",
    ], page)
    page += 1

    # --- Slide 14: Thank You --------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_textbox(slide, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.0),
                "Thank You", size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6),
                "Questions & Live Demo", size=22, color=RGBColor(0xBF, 0xD6, 0xF7))

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Saved presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    sys.exit(build() or 0)
